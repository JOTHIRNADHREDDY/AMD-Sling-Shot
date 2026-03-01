"""
MediKiosk Architecture PPT Generator
Generates a professional PowerPoint presentation covering the full system architecture.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy background
BG_CARD   = RGBColor(0x1A, 0x25, 0x3C)   # Card background
ACCENT    = RGBColor(0x4D, 0xAB, 0xF7)   # Blue accent
ACCENT2   = RGBColor(0x84, 0x5E, 0xF7)   # Purple accent
GREEN     = RGBColor(0x51, 0xCF, 0x66)   # Green
ORANGE    = RGBColor(0xFF, 0x92, 0x2B)   # Orange
RED       = RGBColor(0xFF, 0x6B, 0x6B)   # Red
YELLOW    = RGBColor(0xFF, 0xD4, 0x3B)   # Yellow
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC1, 0xC8, 0xDB)   # Light grey text
TEAL      = RGBColor(0x20, 0xC9, 0x97)   # Teal

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper Functions ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_arrow(slide, start_left, start_top, end_left, end_top, color=ACCENT):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

def card_with_text(slide, left, top, width, height, title, items, card_color=BG_CARD, title_color=ACCENT, border_color=None):
    """Create a card with title and bullet items."""
    shape = add_shape(slide, left, top, width, height, card_color, border_color, Pt(1.5))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = title_color
    p.font.bold = True
    
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = item
        p2.font.size = Pt(10)
        p2.font.color.rgb = LIGHT
        p2.space_before = Pt(2)
    return shape


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_slide_bg(slide, BG_DARK)

# Top accent bar
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "MediKiosk", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Voice-First AI Hospital Self-Service Kiosk", font_size=26, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

# Description
add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(0.8),
             "System Architecture & Process Flow", font_size=20, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Tags bar
tags = ["React 19", "FastAPI", "Llama 3 LLM", "Whisper STT", "Sarvam TTS", "Firebase"]
tag_start = Inches(2.5)
for i, tag in enumerate(tags):
    x = tag_start + Inches(i * 1.45)
    shape = add_shape(slide, x, Inches(5.2), Inches(1.3), Inches(0.38), BG_CARD, ACCENT, Pt(1))
    tf = shape.text_frame
    tf.paragraphs[0].text = tag
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Bottom info
add_text_box(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.5),
             "AMD General Hospital  ·  Telugu · Hindi · Tamil · English  ·  March 2026",
             font_size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Solution
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Problem & Solution", font_size=32, color=WHITE, bold=True)

# Problem card
shape = add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, RED, Pt(1.5))
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16); tf.margin_top = Pt(14)
p = tf.paragraphs[0]; p.text = "The Problem"; p.font.size = Pt(22); p.font.color.rgb = RED; p.font.bold = True

problems = [
    "Long queues at government hospital registration counters",
    "Low-literacy patients unable to use digital forms",
    "Language barriers — patients speak Telugu, Hindi, Tamil",
    "No real-time visibility into queue wait times",
    "Paper-based processes prone to errors & loss",
    "Patients struggle to navigate large hospital buildings"
]
for prob in problems:
    p2 = tf.add_paragraph()
    p2.text = f"✗  {prob}"
    p2.font.size = Pt(13); p2.font.color.rgb = LIGHT; p2.space_before = Pt(10)

# Solution card
shape = add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5), BG_CARD, GREEN, Pt(1.5))
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16); tf.margin_top = Pt(14)
p = tf.paragraphs[0]; p.text = "Our Solution"; p.font.size = Pt(22); p.font.color.rgb = GREEN; p.font.bold = True

solutions = [
    "Voice-first self-service kiosk — just speak naturally",
    "AI understands intent in 4 Indian languages",
    "LLM-powered conversational registration flow",
    "Live queue display with estimated wait times",
    "Digital receipts with QR codes — no paper needed",
    "Step-by-step voice-guided hospital navigation"
]
for sol in solutions:
    p2 = tf.add_paragraph()
    p2.text = f"✓  {sol}"
    p2.font.size = Pt(13); p2.font.color.rgb = LIGHT; p2.space_before = Pt(10)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Technology Stack
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Technology Stack", font_size=32, color=WHITE, bold=True)

# Frontend column
card_with_text(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(5.7), "⚛ Frontend", [
    "React 19 + TypeScript",
    "Vite 6 (build tool)",
    "Tailwind CSS (styling)",
    "Framer Motion (animations)",
    "React Context (state mgmt)",
    "WebSocket (voice streaming)",
    "Web Audio API (recording)",
    "IndexedDB (offline storage)",
    "Firebase SDK (cloud storage)",
    "html5-qrcode (QR scanning)",
    "qrcode.react (QR generation)"
], border_color=ACCENT)

# Backend column
card_with_text(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(5.7), "🐍 Backend", [
    "FastAPI (async Python)",
    "Uvicorn (ASGI server)",
    "SQLAlchemy 2.0 (async ORM)",
    "Alembic (DB migrations)",
    "SQLite / PostgreSQL (database)",
    "Redis (caching — prod)",
    "Firebase Admin SDK (storage)",
    "python-jose (JWT auth)",
    "httpx (async HTTP client)",
    "pydub + FFmpeg (audio)",
    "PyMuPDF (PDF processing)"
], border_color=ACCENT2)

# AI / External column
card_with_text(slide, Inches(9.1), Inches(1.3), Inches(3.8), Inches(5.7), "🤖 AI & External", [
    "Meta Llama 3 8B Instruct (LLM)",
    "OpenAI Whisper Large V3 (STT)",
    "  — via Hugging Face Inference API",
    "",
    "Sarvam AI (Indic TTS)",
    "  — Telugu, Hindi, Tamil, English",
    "",
    "Firebase Cloud Storage",
    "Firebase Analytics",
    "",
    "FFmpeg (audio conversion)"
], border_color=ORANGE)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — High-Level System Architecture
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "System Architecture Overview", font_size=32, color=WHITE, bold=True)

# Frontend box
fe = add_shape(slide, Inches(0.4), Inches(1.3), Inches(4), Inches(5.8), BG_CARD, ACCENT, Pt(2))
tf = fe.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.text = "FRONTEND"; p.font.size = Pt(18); p.font.color.rgb = ACCENT; p.font.bold = True
add_paragraph(tf, "React 19 + TypeScript + Vite", 11, LIGHT, space_before=Pt(8))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "📱 8 Screens", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Language · Home · Registration · Queue", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Navigation · Lab Tests · Complaint · Receipt", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "🎤 Voice Components", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "VoiceOrb · ProcessingIndicator", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "⚙ Service Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "VoiceManager (WebSocket)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "AppBrain (Central Controller)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "IntentRouter (5-Layer Classifier)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "RegistrationFlow (State Machine)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "SessionContext · ClarificationGuard", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "💾 Offline Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "IndexedDB + CloudSyncService", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "KioskContext (Single Source of Truth)", 10, LIGHT, space_before=Pt(2))

# Backend box
be = add_shape(slide, Inches(4.8), Inches(1.3), Inches(4), Inches(5.8), BG_CARD, ACCENT2, Pt(2))
tf = be.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.text = "BACKEND"; p.font.size = Pt(18); p.font.color.rgb = ACCENT2; p.font.bold = True
add_paragraph(tf, "FastAPI + Python (Async)", 11, LIGHT, space_before=Pt(8))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "🌐 API Layer (v1)", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "/voice/stream (WebSocket)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "/registration · /queue · /map", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "/lab-tests · /storage (16 endpoints)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "🧠 AI / Services Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "LLM Orchestrator (pipeline controller)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "LLM Tools (register/queue/directions)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Command Engine (safety validator)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Screen Capabilities (whitelists)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Conversation Memory (10-turn)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Workflow State Machine", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "💽 Data Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "SQLAlchemy (async) → SQLite/PostgreSQL", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Firebase Admin SDK → Cloud Storage", 10, LIGHT, space_before=Pt(2))

# External services box
ext = add_shape(slide, Inches(9.2), Inches(1.3), Inches(3.7), Inches(5.8), BG_CARD, ORANGE, Pt(2))
tf = ext.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.text = "EXTERNAL SERVICES"; p.font.size = Pt(18); p.font.color.rgb = ORANGE; p.font.bold = True
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "🗣 Speech-to-Text", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "OpenAI Whisper Large V3", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "via Hugging Face Inference API", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "🤖 Large Language Model", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Meta Llama 3 8B Instruct", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "via Hugging Face Inference API", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "🔊 Text-to-Speech", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Sarvam AI (Indic languages)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Telugu · Hindi · Tamil · English", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "☁ Cloud Storage", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Firebase Cloud Storage", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Signed URLs + Analytics", 10, LIGHT, space_before=Pt(2))

# Connection arrows (labels)
add_text_box(slide, Inches(4.05), Inches(3.8), Inches(0.8), Inches(0.4),
             "◄──►", font_size=18, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.8), Inches(4.15), Inches(1.3), Inches(0.35),
             "WebSocket + REST", font_size=9, color=YELLOW, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(8.5), Inches(3.8), Inches(0.8), Inches(0.4),
             "◄──►", font_size=18, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.3), Inches(4.15), Inches(1.2), Inches(0.35),
             "HTTPS APIs", font_size=9, color=YELLOW, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Frontend Architecture Detail
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Frontend Architecture — React 19 SPA", font_size=32, color=WHITE, bold=True)

# Screens section
screens_data = [
    ("Language", "Language gate"), ("Home", "6 touch tiles"),
    ("Registration", "Multi-step wizard"), ("Queue", "Live status"),
    ("Navigation", "Wayfinding"), ("Lab Tests", "Results viewer"),
    ("Complaint", "Filing form"), ("Receipt", "QR + Token")
]
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(3), Inches(0.4),
             "📱 Screens (8)", font_size=16, color=ACCENT, bold=True)
for i, (name, desc) in enumerate(screens_data):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 1.9)
    y = Inches(1.7 + row * 0.55)
    s = add_shape(slide, x, y, Inches(1.75), Inches(0.45), BG_CARD, ACCENT, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = f"{name}"; p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(8); p2.font.color.rgb = LIGHT

# Components section
add_text_box(slide, Inches(0.5), Inches(4.0), Inches(3), Inches(0.4),
             "🧩 Components", font_size=16, color=ACCENT2, bold=True)
components = ["VoiceOrb — voice I/O button", "AdaptiveHeader — smart nav bar",
              "ProcessingIndicator — LLM wait", "HelpModal — context help"]
for i, comp in enumerate(components):
    s = add_shape(slide, Inches(0.5), Inches(4.45 + i * 0.48), Inches(3.5), Inches(0.4), BG_CARD, ACCENT2, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(4)
    p = tf.paragraphs[0]; p.text = comp; p.font.size = Pt(10); p.font.color.rgb = LIGHT

# Services section
add_text_box(slide, Inches(4.5), Inches(1.2), Inches(4.5), Inches(0.4),
             "⚙ Service Layer", font_size=16, color=GREEN, bold=True)
services = [
    ("VoiceManager", "WebSocket + MediaRecorder + TTS playback"),
    ("AppBrain", "Central controller — all UI actions flow through here"),
    ("IntentRouter", "5-layer classifier: exact → numeric → keyword → medical → LLM"),
    ("VoiceCommandEngine", "Pipeline: transcript → filter → intent → dispatch"),
    ("RegistrationFlow", "State machine: IDLE→MOBILE→NAME→AGE→GENDER→DEPT→CONFIRM"),
    ("SessionContext", "Cross-flow memory in sessionStorage + FlowLock"),
    ("ClarificationGuard", "Prevents LLM infinite clarification loops"),
    ("VoiceNormalizer", "Entity extraction: token, name, age, gender, mobile"),
    ("api.ts", "Typed REST client for all backend endpoints"),
    ("offlineStorage", "IndexedDB — 5 stores for offline-first support"),
]
for i, (name, desc) in enumerate(services):
    y = Inches(1.7 + i * 0.5)
    s = add_shape(slide, Inches(4.5), y, Inches(4.5), Inches(0.44), BG_CARD, GREEN, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(10); p.font.color.rgb = GREEN; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(8); p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# State management
add_text_box(slide, Inches(9.5), Inches(1.2), Inches(3.5), Inches(0.4),
             "🏪 KioskContext (State)", font_size=16, color=YELLOW, bold=True)
state_items = [
    "Screen navigation state",
    "Voice pipeline state (IDLE/LISTENING/PROCESSING/SPEAKING)",
    "Patient data & registration result",
    "Queue data (auto-poll 30s)",
    "Map directions",
    "Registration flow state machine",
    "Cloud storage state (upload progress)",
    "Workflow state (IDLE→COLLECTING→CONFIRM→SUBMITTED)",
    "Interaction lock (prevents double-actions)",
    "Language & translations"
]
card_with_text(slide, Inches(9.5), Inches(1.7), Inches(3.5), Inches(5.3),
               "Single Source of Truth", state_items, BG_CARD, YELLOW)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Backend Architecture Detail
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Backend Architecture — FastAPI (Async Python)", font_size=32, color=WHITE, bold=True)

# API Layer
add_text_box(slide, Inches(0.4), Inches(1.15), Inches(4), Inches(0.4),
             "🌐 API Layer (v1)", font_size=16, color=ACCENT, bold=True)
api_endpoints = [
    ("/voice/stream", "WebSocket", "Real-time audio streaming + LLM response"),
    ("/voice/intent", "POST", "REST-based intent classification"),
    ("/voice/chat", "POST", "Generic LLM chat endpoint"),
    ("/registration/register", "POST", "Create patient + queue token"),
    ("/registration/lookup/{token}", "GET", "QR/token lookup"),
    ("/queue/status", "GET", "Live department queue aggregation"),
    ("/map/directions", "GET", "Dijkstra-based wayfinding"),
    ("/lab-tests/*", "CRUD", "Scan upload, results, pagination"),
    ("/storage/*", "16 endpoints", "Upload/download/copy/move/delete files"),
]
for i, (path, method, desc) in enumerate(api_endpoints):
    y = Inches(1.55 + i * 0.47)
    s = add_shape(slide, Inches(0.4), y, Inches(4.2), Inches(0.42), BG_CARD, ACCENT, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = f"{method}  {path}"; p.font.size = Pt(9); p.font.color.rgb = ACCENT; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(8); p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# AI Services Layer
add_text_box(slide, Inches(5), Inches(1.15), Inches(4.2), Inches(0.4),
             "🧠 AI / Services Layer", font_size=16, color=ACCENT2, bold=True)
ai_services = [
    ("LLM Orchestrator", "Core pipeline: Audio → STT → Prompt → LLM → Tool → TTS"),
    ("LLM Tools", "register_patient, get_queue, get_directions, lookup_token, submit_complaint"),
    ("Command Engine", "Translates ToolResult → validated CommandEnvelope. Safety layer."),
    ("Screen Capabilities", "Per-screen action whitelists + field-level edit control"),
    ("Conversation Memory", "Session-scoped: 10-turn history, partial data, clarification ctx"),
    ("Workflow State", "State machine: IDLE→COLLECTING→CONFIRMATION→SUBMITTED→COMPLETE"),
    ("JSON Validator", "Validates LLM output JSON structure before execution"),
    ("Fallback Intent", "Keyword-based fallback when LLM is unreachable"),
]
for i, (name, desc) in enumerate(ai_services):
    y = Inches(1.55 + i * 0.52)
    s = add_shape(slide, Inches(5), y, Inches(4.2), Inches(0.46), BG_CARD, ACCENT2, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(10); p.font.color.rgb = ACCENT2; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(8); p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# Data Layer
add_text_box(slide, Inches(9.6), Inches(1.15), Inches(3.5), Inches(0.4),
             "💽 Data Layer", font_size=16, color=ORANGE, bold=True)
data_items = [
    ("SQLAlchemy 2.0 (async)", "aiosqlite dev / PostgreSQL prod"),
    ("Domain Models", "User, OPRegistration, QueuePosition, LabTestScan, HospitalBranch"),
    ("Alembic Migrations", "Async migration support"),
    ("Firebase Admin SDK", "Signed URLs, upload/download, retry logic"),
    ("Redis (prod)", "Session caching, rate limiting"),
]
for i, (name, desc) in enumerate(data_items):
    y = Inches(1.55 + i * 0.65)
    s = add_shape(slide, Inches(9.6), y, Inches(3.4), Inches(0.58), BG_CARD, ORANGE, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(10); p.font.color.rgb = ORANGE; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(8); p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# Config
add_text_box(slide, Inches(9.6), Inches(5.0), Inches(3.5), Inches(0.35),
             "⚙ Configuration (pydantic-settings)", font_size=12, color=YELLOW, bold=True)
config_items = ["DB creds · Redis URL · JWT settings", "HF_TOKEN · SARVAM_API_KEY · FFmpeg path",
                "Rate limiting: 30 req/min, 10MB max, 30s timeout"]
card_with_text(slide, Inches(9.6), Inches(5.35), Inches(3.4), Inches(1.5),
               "", config_items, BG_CARD, YELLOW)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Voice Command Pipeline (End-to-End)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Voice Command Pipeline — End-to-End Flow", font_size=32, color=WHITE, bold=True)

# Pipeline steps
steps = [
    ("1", "Patient Speaks", "Taps VoiceOrb or\nauto-trigger", ACCENT, Inches(0.3)),
    ("2", "Audio Capture", "MediaRecorder →\nWebSocket binary", ACCENT, Inches(1.65)),
    ("3", "Whisper STT", "Speech-to-Text\n(HF Inference API)", ACCENT2, Inches(3.0)),
    ("4", "LLM Processing", "Llama 3 8B Instruct\n+ system prompt", ACCENT2, Inches(4.35)),
    ("5", "Tool Execution", "register · queue\ndirections · lookup", ORANGE, Inches(5.7)),
    ("6", "Command Engine", "Validate against\nScreen Capabilities", RED, Inches(7.05)),
    ("7", "TTS Generation", "Sarvam AI → Indic\naudio response", GREEN, Inches(8.4)),
    ("8", "UI Update", "Navigate · Fill form\nClick · Show data", GREEN, Inches(9.75)),
    ("9", "Speak Response", "VoiceOrb plays\nTTS audio to patient", TEAL, Inches(11.1)),
]

# Draw step boxes with arrows
for i, (num, title, desc, color, x_pos) in enumerate(steps):
    y = Inches(1.6)
    s = add_shape(slide, x_pos, y, Inches(1.25), Inches(1.8), BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(8)
    
    # Number circle
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(20); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = title; p2.font.size = Pt(10); p2.font.color.rgb = WHITE; p2.font.bold = True; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
    p3 = tf.add_paragraph(); p3.text = desc; p3.font.size = Pt(8); p3.font.color.rgb = LIGHT; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(4)
    
    # Arrow
    if i < len(steps) - 1:
        ax = x_pos + Inches(1.3)
        add_text_box(slide, ax, Inches(2.2), Inches(0.4), Inches(0.4), "→", font_size=20, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Parallel local path
add_shape(slide, Inches(0.3), Inches(3.8), Inches(12.5), Inches(0.04), YELLOW)
add_text_box(slide, Inches(0.3), Inches(3.95), Inches(12.5), Inches(0.35),
             "⚡ PARALLEL LOCAL PATH: Transcript also goes to IntentRouter (5-layer classifier) for instant response without LLM round-trip",
             font_size=11, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Safety pipeline detail
add_text_box(slide, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
             "🔒 LLM Safety Pipeline — Raw LLM output NEVER reaches the UI directly", font_size=16, color=RED, bold=True)

safety_steps = [
    ("LLM Raw Output", RED, "⚠ Untrusted"),
    ("JSON Validator", YELLOW, "Structure check"),
    ("Command Engine", ORANGE, "Action validation"),
    ("Screen Capabilities", ACCENT2, "Whitelist check"),
    ("Field Whitelist", ACCENT, "Per-form validation"),
    ("CommandEnvelope", GREEN, "✓ Safe for UI"),
]
for i, (name, color, desc) in enumerate(safety_steps):
    x = Inches(0.5 + i * 2.15)
    s = add_shape(slide, x, Inches(5.15), Inches(1.95), Inches(0.95), BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(8); tf.margin_top = Pt(6)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(11); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(9); p2.font.color.rgb = LIGHT; p2.alignment = PP_ALIGN.CENTER
    
    if i < len(safety_steps) - 1:
        ax = x + Inches(2.0)
        add_text_box(slide, ax, Inches(5.35), Inches(0.3), Inches(0.35), "→", font_size=18, color=YELLOW, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
             "Every voice action passes through 5 validation stages before it can modify any UI element",
             font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Registration Process Flow
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Patient Registration Flow", font_size=32, color=WHITE, bold=True)

# Flow steps
reg_steps = [
    ("1", "Arrive", "Patient walks\nto kiosk", ACCENT),
    ("2", "Language", "Select: Telugu\nHindi · Tamil · English", ACCENT),
    ("3", "Home", "Tap 'Register'\ntile or say it", ACCENT2),
    ("4", "Mobile", "Enter 10-digit\nnumber (pad/voice)", GREEN),
    ("5", "Details", "Name + Age +\nGender (voice/touch)", GREEN),
    ("6", "Department", "Select from\nemoji grid or voice", GREEN),
    ("7", "Confirm", "Review all\ndetails", YELLOW),
    ("8", "Submit", "Backend creates\nOPRegistration + Token", ORANGE),
    ("9", "Receipt", "QR code + Token\ne.g. ENT-007", TEAL),
]

for i, (num, title, desc, color) in enumerate(reg_steps):
    x = Inches(0.3 + i * 1.42)
    y = Inches(1.5)
    s = add_shape(slide, x, y, Inches(1.3), Inches(1.8), BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = f"Step {num}"; p.font.size = Pt(10); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = title; p2.font.size = Pt(13); p2.font.color.rgb = WHITE; p2.font.bold = True; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
    p3 = tf.add_paragraph(); p3.text = desc; p3.font.size = Pt(9); p3.font.color.rgb = LIGHT; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(6)
    
    if i < len(reg_steps) - 1:
        ax = x + Inches(1.33)
        add_text_box(slide, ax, Inches(2.1), Inches(0.15), Inches(0.4), "→", font_size=16, color=YELLOW, alignment=PP_ALIGN.CENTER)

# State machine
add_text_box(slide, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4),
             "Registration State Machine (RegistrationFlow.ts)", font_size=16, color=GREEN, bold=True)

states = ["IDLE", "MOBILE", "NAME", "AGE", "GENDER", "DEPARTMENT", "CONFIRM", "SUBMITTED"]
state_colors = [LIGHT, ACCENT, ACCENT, ACCENT2, ACCENT2, GREEN, YELLOW, TEAL]
for i, (state, col) in enumerate(zip(states, state_colors)):
    x = Inches(0.5 + i * 1.55)
    s = add_shape(slide, x, Inches(4.2), Inches(1.35), Inches(0.55), BG_CARD, col, Pt(2))
    tf = s.text_frame; tf.margin_top = Pt(6)
    p = tf.paragraphs[0]; p.text = state; p.font.size = Pt(12); p.font.color.rgb = col; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    if i < len(states) - 1:
        add_text_box(slide, x + Inches(1.38), Inches(4.3), Inches(0.2), Inches(0.35), "→", font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Features
add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.35),
             "Key Features:", font_size=14, color=WHITE, bold=True)
features = [
    "FlowLock — blocks unrelated navigation during registration",
    "Auto-advances past already-filled fields",
    "LLM conversationally guides each step in patient's language",
    "Backend creates OPRegistration + QueuePosition atomically",
    "Receipt uploaded to Firebase + saved to IndexedDB offline"
]
for i, feat in enumerate(features):
    add_text_box(slide, Inches(0.5), Inches(5.6 + i * 0.32), Inches(12), Inches(0.35),
                 f"  ✓  {feat}", font_size=11, color=LIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Complete Application Flow
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Complete Application Process Flow", font_size=32, color=WHITE, bold=True)

# Home screen center
home = add_shape(slide, Inches(5.2), Inches(1.2), Inches(3), Inches(0.8), BG_CARD, ACCENT2, Pt(2))
tf = home.text_frame; tf.margin_top = Pt(6)
p = tf.paragraphs[0]; p.text = "🏠 HOME SCREEN"; p.font.size = Pt(16); p.font.color.rgb = ACCENT2; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "6 Touch Tiles + Voice Orb"; p2.font.size = Pt(10); p2.font.color.rgb = LIGHT; p2.alignment = PP_ALIGN.CENTER

# 6 service flows radiating out
services_flow = [
    ("🏥 OP Registration", "Voice/touch wizard → DB\n→ Token + QR Receipt", ACCENT, Inches(0.3), Inches(2.5)),
    ("📋 Queue Status", "GET /queue/status → DB\nAuto-poll 30s · Live ETA", GREEN, Inches(0.3), Inches(4.3)),
    ("🗺 Navigation", "GET /map/directions\nStep-by-step wayfinding", TEAL, Inches(0.3), Inches(6.1)),
    ("🧪 Lab Tests", "CRUD /lab-tests\nUpload scans → Firebase", ORANGE, Inches(8.8), Inches(2.5)),
    ("📝 Complaints", "Workflow: COLLECTING →\nCONFIRM → SUBMITTED", RED, Inches(8.8), Inches(4.3)),
    ("🔍 Token Lookup", "GET /registration/lookup\nQR scan or voice input", YELLOW, Inches(8.8), Inches(6.1)),
]

for title, desc, color, x, y in services_flow:
    s = add_shape(slide, x, y, Inches(4.2), Inches(1.4), BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(14); p.font.color.rgb = color; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(10); p2.font.color.rgb = LIGHT; p2.space_before = Pt(4)

# Labels from home to services
add_text_box(slide, Inches(4.5), Inches(2.8), Inches(0.8), Inches(0.3), "◄──", font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(4.6), Inches(0.8), Inches(0.3), "◄──", font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(6.4), Inches(0.8), Inches(0.3), "◄──", font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.1), Inches(2.8), Inches(0.8), Inches(0.3), "──►", font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.1), Inches(4.6), Inches(0.8), Inches(0.3), "──►", font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.1), Inches(6.4), Inches(0.8), Inches(0.3), "──►", font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Voice overlay note
vo = add_shape(slide, Inches(4.8), Inches(2.3), Inches(3.5), Inches(0.45), BG_CARD, GREEN, Pt(1.5))
tf = vo.text_frame; tf.margin_top = Pt(4); tf.margin_left = Pt(8)
p = tf.paragraphs[0]; p.text = "🎤 VoiceOrb available on EVERY screen"; p.font.size = Pt(11); p.font.color.rgb = GREEN; p.font.bold = True; p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Design Patterns & Safety
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Design Patterns & Safety Architecture", font_size=32, color=WHITE, bold=True)

patterns = [
    ("Command Pattern", "All UI mutations are wrapped in typed CommandEnvelope objects", ACCENT),
    ("State Machine", "RegistrationFlow + WorkflowState enforce legal transitions only", GREEN),
    ("Mediator / Orchestrator", "LLM Orchestrator coordinates STT → LLM → Tool → TTS pipeline", ACCENT2),
    ("Observer / Pub-Sub", "VoiceManager.subscribe() for decoupled event notification", TEAL),
    ("Strategy (Layered)", "IntentRouter 5-layer detection — each layer has different matching strategy", ORANGE),
    ("Capability Matrix", "Per-screen action whitelists + field-level edit control", YELLOW),
    ("Offline-First", "IndexedDB fallback when cloud unreachable; auto-sync when restored", ACCENT),
    ("Singleton", "voiceManager, flowLock, sessionContext — global cross-component instances", LIGHT),
    ("Guard Pattern", "ClarificationGuard + FlowLock + confidence filter prevent loops & drift", RED),
    ("Interaction Lock", "400ms lock + 350ms debounce prevents double-execution of actions", ACCENT2),
]

# Left column
for i in range(5):
    name, desc, color = patterns[i]
    y = Inches(1.3 + i * 1.1)
    s = add_shape(slide, Inches(0.4), y, Inches(6.2), Inches(0.95), BG_CARD, color, Pt(1.5))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(14); p.font.color.rgb = color; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = LIGHT; p2.space_before = Pt(4)

# Right column
for i in range(5, 10):
    name, desc, color = patterns[i]
    y = Inches(1.3 + (i - 5) * 1.1)
    s = add_shape(slide, Inches(6.9), y, Inches(6.2), Inches(0.95), BG_CARD, color, Pt(1.5))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(14); p.font.color.rgb = color; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(11); p2.font.color.rgb = LIGHT; p2.space_before = Pt(4)

# Safety highlight
add_text_box(slide, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
             "🔒 Core Safety Principle: LLM interprets intent but NEVER directly mutates state — all actions pass through validated command envelopes, "
             "screen capability checks, field whitelists, and state machines before reaching the UI.",
             font_size=11, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Offline & Storage Architecture
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Offline-First Storage & Data Flow", font_size=32, color=WHITE, bold=True)

# User action
s = add_shape(slide, Inches(4.5), Inches(1.3), Inches(4.5), Inches(0.7), BG_CARD, ACCENT, Pt(2))
tf = s.text_frame; tf.margin_top = Pt(8)
p = tf.paragraphs[0]; p.text = "User Action: Save receipt / Upload scan / Store document"; p.font.size = Pt(12); p.font.color.rgb = ACCENT; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# Decision
add_text_box(slide, Inches(5.3), Inches(2.05), Inches(3), Inches(0.4), "▼  Cloud available?", font_size=12, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Cloud path (left)
s = add_shape(slide, Inches(1.5), Inches(2.7), Inches(4), Inches(0.7), BG_CARD, GREEN, Pt(2))
tf = s.text_frame; tf.margin_top = Pt(6)
p = tf.paragraphs[0]; p.text = "✓ YES → Firebase Cloud Storage"; p.font.size = Pt(13); p.font.color.rgb = GREEN; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "Upload with signed URL → Store metadata in DB"; p2.font.size = Pt(10); p2.font.color.rgb = LIGHT; p2.alignment = PP_ALIGN.CENTER

# Offline path (right)
s = add_shape(slide, Inches(7.8), Inches(2.7), Inches(4), Inches(0.7), BG_CARD, RED, Pt(2))
tf = s.text_frame; tf.margin_top = Pt(6)
p = tf.paragraphs[0]; p.text = "✗ NO → IndexedDB Offline Store"; p.font.size = Pt(13); p.font.color.rgb = RED; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "5 stores: labScans · registrations · receipts · docs · reports"; p2.font.size = Pt(10); p2.font.color.rgb = LIGHT; p2.alignment = PP_ALIGN.CENTER

# Sync arrow
s = add_shape(slide, Inches(7.8), Inches(3.7), Inches(4), Inches(0.55), BG_CARD, YELLOW, Pt(1.5))
tf = s.text_frame; tf.margin_top = Pt(6)
p = tf.paragraphs[0]; p.text = "⏳ cloudSyncService → auto-sync when network restored"; p.font.size = Pt(10); p.font.color.rgb = YELLOW; p.alignment = PP_ALIGN.CENTER

# Database model
add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.4),
             "Database Models (SQLAlchemy)", font_size=18, color=ACCENT2, bold=True)

models = [
    ("User", "id · name · email · role\nhashed_password · branch_id", ACCENT),
    ("OPRegistration", "id · patient_name · age · gender\nmobile · department · token\nform_data (JSON) · status", GREEN),
    ("QueuePosition", "id · registration_id · department\nposition · status · called_at\nestimated_wait_minutes", ORANGE),
    ("LabTestScan", "id · patient_name · scan_type\nfile_url · status · results\nuploaded_at · processed_at", ACCENT2),
    ("HospitalBranch", "id · name · code\naddress · is_active", TEAL),
]

for i, (name, fields, color) in enumerate(models):
    x = Inches(0.3 + i * 2.6)
    s = add_shape(slide, x, Inches(5.3), Inches(2.4), Inches(1.8), BG_CARD, color, Pt(1.5))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(13); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = fields; p2.font.size = Pt(9); p2.font.color.rgb = LIGHT; p2.space_before = Pt(6); p2.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Multilingual Support
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Multilingual & Accessibility Architecture", font_size=32, color=WHITE, bold=True)

languages = [
    ("తెలుగు", "Telugu", "Primary language of Andhra Pradesh", ACCENT),
    ("हिन्दी", "Hindi", "National language — wide reach", ACCENT2),
    ("தமிழ்", "Tamil", "For Tamil-speaking patients", GREEN),
    ("English", "English", "Fallback / literate patients", ORANGE),
]

for i, (script, name, desc, color) in enumerate(languages):
    x = Inches(0.5 + i * 3.2)
    s = add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(1.5), BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(12); tf.margin_top = Pt(12)
    p = tf.paragraphs[0]; p.text = script; p.font.size = Pt(28); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = name; p2.font.size = Pt(14); p2.font.color.rgb = WHITE; p2.font.bold = True; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
    p3 = tf.add_paragraph(); p3.text = desc; p3.font.size = Pt(10); p3.font.color.rgb = LIGHT; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(4)

# How multilingual works
add_text_box(slide, Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
             "How Multilingual Pipeline Works", font_size=18, color=YELLOW, bold=True)

ml_steps = [
    ("Language Gate", "User selects language\non first screen", ACCENT),
    ("STT (Whisper)", "Transcribes audio\nin detected language", ACCENT2),
    ("LLM Prompt", "System prompt instructs\nLlama 3 to respond in\ntarget language script", GREEN),
    ("Localized Messages", "_LOCALIZED_MESSAGES dict\nin orchestrator — 4 language\ntemplates for all responses", ORANGE),
    ("TTS (Sarvam AI)", "Generates speech audio\nin patient's language\nIndic language support", TEAL),
]

for i, (title, desc, color) in enumerate(ml_steps):
    x = Inches(0.3 + i * 2.6)
    s = add_shape(slide, x, Inches(3.7), Inches(2.4), Inches(1.8), BG_CARD, color, Pt(1.5))
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Pt(10); tf.margin_top = Pt(10)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(13); p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(10); p2.font.color.rgb = LIGHT; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    
    if i < len(ml_steps) - 1:
        add_text_box(slide, x + Inches(2.43), Inches(4.3), Inches(0.25), Inches(0.35), "→", font_size=16, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Accessibility features
add_text_box(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(0.35),
             "Accessibility Features for Low-Literacy Patients", font_size=16, color=GREEN, bold=True)
access_features = [
    "Large touch tiles with emoji icons — no reading required",
    "Voice-first interaction — speak naturally in any supported language",
    "VoiceOrb gives visual feedback: color-coded states (listening/processing/speaking)",
    "LLM conversationally guides step-by-step — never shows a blank form",
    "Auto-fill via voice — patient doesn't need to type anything",
]
for i, feat in enumerate(access_features):
    add_text_box(slide, Inches(0.5), Inches(6.3 + i * 0.22), Inches(12), Inches(0.25),
                 f"  ✓  {feat}", font_size=10, color=LIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Summary / Thank You
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "MediKiosk", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.6),
             "Voice-First AI Hospital Self-Service Kiosk", font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

summary_items = [
    "🎤  Voice-first design for low-literacy patients in 4 Indian languages",
    "🤖  LLM-orchestrated (Llama 3) with Whisper STT and Sarvam TTS",
    "🔒  5-layer safety pipeline — LLM output validated before reaching UI",
    "⚡  Dual processing: instant local intent + full LLM pipeline in parallel",
    "📱  8 screens: Registration · Queue · Navigation · Lab Tests · Complaints",
    "💾  Offline-first with IndexedDB fallback + Firebase cloud sync",
    "🏥  Deterministic state machines for registration & workflow flows",
]

for i, item in enumerate(summary_items):
    y = Inches(3.2 + i * 0.45)
    add_text_box(slide, Inches(2), y, Inches(9), Inches(0.4),
                 item, font_size=15, color=LIGHT, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "AMD General Hospital  ·  March 2026", font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Bottom accent bar
add_rect(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), ACCENT)


# ── Save PPT ──
output_path = os.path.join(r"e:\AMD\medikisok", "MediKiosk_Architecture.pptx")
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
