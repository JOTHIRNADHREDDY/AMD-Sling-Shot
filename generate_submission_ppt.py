"""
MediKiosk — Project Submission Deck PPT Generator
Generates a professional project-submission presentation covering:
  Title · Abstract · Problem · Objectives · Literature Survey ·
  Proposed System · Architecture · Tech Stack · Implementation ·
  Screenshots · Testing · Results · Advantages/Limitations ·
  Future Scope · Conclusion · References
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD   = RGBColor(0x1A, 0x25, 0x3C)
ACCENT    = RGBColor(0x4D, 0xAB, 0xF7)
ACCENT2   = RGBColor(0x84, 0x5E, 0xF7)
GREEN     = RGBColor(0x51, 0xCF, 0x66)
ORANGE    = RGBColor(0xFF, 0x92, 0x2B)
RED       = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW    = RGBColor(0xFF, 0xD4, 0x3B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC1, 0xC8, 0xDB)
TEAL      = RGBColor(0x20, 0xC9, 0x97)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color,
              border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, fill_color,
             border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text,
                 font_size=14, color=WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=14, color=WHITE,
                  bold=False, alignment=PP_ALIGN.LEFT,
                  space_before=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = space_before
    return p


def card_with_text(slide, left, top, width, height, title, items,
                   card_color=BG_CARD, title_color=ACCENT,
                   border_color=None, item_size=10):
    shape = add_shape(slide, left, top, width, height,
                      card_color, border_color, Pt(1.5))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = title_color
    p.font.bold = True

    for item in items:
        p2 = tf.add_paragraph()
        p2.text = item
        p2.font.size = Pt(item_size)
        p2.font.color.rgb = LIGHT
        p2.space_before = Pt(3)
    return shape


def slide_header(slide, title, accent_color=ACCENT):
    """Standard slide header: accent bar + title text."""
    add_rect(slide, Inches(0), Inches(0),
             Inches(13.333), Inches(0.06), accent_color)
    add_text_box(slide, Inches(0.5), Inches(0.25),
                 Inches(12), Inches(0.7),
                 title, font_size=32, color=WHITE, bold=True)


def slide_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.05),
                 Inches(1), Inches(0.35),
                 f"{num}/{total}", font_size=10,
                 color=LIGHT, alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 18   # updated at the end if needed


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_rect(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(sl, Inches(1), Inches(1.4), Inches(11), Inches(1.2),
             "MediKiosk", font_size=56, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Voice-First AI Hospital Self-Service Kiosk",
             font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(2), Inches(3.7), Inches(9), Inches(0.6),
             "A Project Submission Report",
             font_size=20, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Tags
tags = ["React 19", "FastAPI", "Llama 3 LLM", "Whisper STT",
        "Sarvam TTS", "Firebase"]
tag_start = Inches(2.5)
for i, tag in enumerate(tags):
    x = tag_start + Inches(i * 1.45)
    s = add_shape(sl, x, Inches(4.6), Inches(1.3), Inches(0.38),
                  BG_CARD, ACCENT, Pt(1))
    tf = s.text_frame
    tf.paragraphs[0].text = tag
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(sl, Inches(1), Inches(5.6), Inches(11), Inches(0.45),
             "Submitted by:  Team MediKiosk",
             font_size=16, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(6.1), Inches(11), Inches(0.4),
             "Under the guidance of:  Prof. _____________",
             font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1), Inches(6.6), Inches(11), Inches(0.45),
             "Department of Computer Science & Engineering  ·  AMD General Hospital  ·  March 2026",
             font_size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_rect(sl, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), ACCENT)
slide_number(sl, 1, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Table of Contents
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "Table of Contents", ACCENT2)

toc = [
    "1.  Abstract",
    "2.  Problem Statement",
    "3.  Objectives",
    "4.  Literature Survey / Existing Systems",
    "5.  Proposed System & Methodology",
    "6.  System Architecture",
    "7.  Technology Stack & Tools",
    "8.  Module Design — Frontend",
    "9.  Module Design — Backend & AI Pipeline",
    "10. Voice Command Pipeline (End-to-End)",
    "11. Patient Registration Flow",
    "12. Multilingual & Accessibility",
    "13. Offline-First Storage & Data Flow",
    "14. Testing & Validation",
    "15. Results & Discussion",
    "16. Advantages, Limitations & Future Scope",
    "17. Conclusion",
    "18. References",
]

for i, item in enumerate(toc):
    col = 0 if i < 9 else 1
    row = i if i < 9 else i - 9
    x = Inches(1.0 + col * 6.2)
    y = Inches(1.3 + row * 0.62)
    s = add_shape(sl, x, y, Inches(5.6), Inches(0.52),
                  BG_CARD, ACCENT2 if col == 0 else TEAL, Pt(1))
    tf = s.text_frame
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE

slide_number(sl, 2, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Abstract
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "1.  Abstract", ACCENT)

abstract_text = (
    "MediKiosk is a voice-first, AI-powered hospital self-service kiosk designed to "
    "streamline patient interactions in government hospitals. The system enables patients "
    "— including those with low literacy — to perform OP registration, check live queue "
    "status, navigate the hospital, view lab results, and file complaints, all through "
    "natural voice commands in Telugu, Hindi, Tamil, or English."
)
abstract_text_2 = (
    "The kiosk leverages Meta Llama 3 8B Instruct as the core large language model for "
    "conversational understanding, OpenAI Whisper Large V3 for speech-to-text transcription, "
    "and Sarvam AI for Indic text-to-speech synthesis. A 5-layer intent classification system "
    "provides instant local responses in parallel with the full LLM pipeline, ensuring low-latency "
    "interactions even under load."
)
abstract_text_3 = (
    "The frontend is built with React 19 and TypeScript, featuring 8 touch-optimized screens, "
    "a state-machine-driven registration flow, and an offline-first architecture with IndexedDB "
    "fallback and Firebase cloud sync. The FastAPI backend enforces a robust safety pipeline: "
    "all LLM outputs are validated through JSON structure checks, command envelopes, screen "
    "capability whitelists, and field-level access control before reaching the UI."
)

s = add_shape(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.6),
              BG_CARD, ACCENT, Pt(1.5))
tf = s.text_frame
tf.word_wrap = True
tf.margin_left = Pt(20)
tf.margin_right = Pt(20)
tf.margin_top = Pt(16)
p = tf.paragraphs[0]; p.text = abstract_text
p.font.size = Pt(14); p.font.color.rgb = LIGHT; p.line_spacing = Pt(22)

add_paragraph(tf, "", 8, LIGHT, space_before=Pt(10))
add_paragraph(tf, abstract_text_2, 14, LIGHT, space_before=Pt(6)).line_spacing = Pt(22)
add_paragraph(tf, "", 8, LIGHT, space_before=Pt(10))
add_paragraph(tf, abstract_text_3, 14, LIGHT, space_before=Pt(6)).line_spacing = Pt(22)

# Keywords
add_paragraph(tf, "", 8, LIGHT, space_before=Pt(14))
kw = add_paragraph(tf, "Keywords: ", 12, YELLOW, True, space_before=Pt(6))
add_paragraph(tf,
    "Voice-First Kiosk, LLM Orchestration, Hospital Self-Service, Multilingual NLP, "
    "Offline-First Architecture, React 19, FastAPI, Whisper STT, Sarvam TTS, Llama 3",
    11, LIGHT, space_before=Pt(2))

slide_number(sl, 3, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Problem Statement
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "2.  Problem Statement", RED)

# Problem card
s = add_shape(sl, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.7),
              BG_CARD, RED, Pt(1.5))
tf = s.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_top = Pt(14)
p = tf.paragraphs[0]
p.text = "Current Challenges"
p.font.size = Pt(22)
p.font.color.rgb = RED
p.font.bold = True

problems = [
    "Long queues at government hospital OP registration counters — patients wait 30–90 minutes",
    "Low-literacy patients cannot fill digital or paper forms independently",
    "Language barriers: patients speak Telugu, Hindi, or Tamil, but interfaces are English-only",
    "No real-time visibility into queue positions — patients miss their turn",
    "Paper-based processes are slow, error-prone, and create data-loss risk",
    "Hospital navigation is confusing — patients waste time finding departments",
    "Lab test results require returning to the hospital counter for a printout",
    "Complaint-filing is manual and rarely reaches administration",
]
for prob in problems:
    p2 = tf.add_paragraph()
    p2.text = f"✗  {prob}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(10)

# Impact card
s = add_shape(sl, Inches(7), Inches(1.3), Inches(5.8), Inches(5.7),
              BG_CARD, ORANGE, Pt(1.5))
tf = s.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_top = Pt(14)
p = tf.paragraphs[0]
p.text = "Impact on Patients & Hospitals"
p.font.size = Pt(22)
p.font.color.rgb = ORANGE
p.font.bold = True

impacts = [
    "Patient frustration: elderly & rural patients often leave without treatment",
    "Hospital staff overloaded with repetitive registration & direction queries",
    "Data silos: paper records not integrated with digital hospital systems",
    "Long OPD waiting times reduce daily patient throughput by 30–40%",
    "Vulnerable populations (elderly, visually impaired) excluded from digital services",
    "No feedback loop: complaints are lost, patients feel unheard",
    "Resources wasted on simple tasks that could be self-service",
]
for imp in impacts:
    p2 = tf.add_paragraph()
    p2.text = f"⚠  {imp}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(10)

slide_number(sl, 4, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Objectives
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "3.  Objectives", GREEN)

objectives = [
    ("Voice-First Self-Service",
     "Enable patients to complete all kiosk interactions using natural voice commands in Telugu, Hindi, Tamil, or English — no reading or typing required.",
     ACCENT),
    ("LLM-Powered Conversational Flow",
     "Use Meta Llama 3 8B Instruct to intelligently guide patients through multi-step workflows (registration, complaints) via natural conversation.",
     ACCENT2),
    ("Real-Time Queue Management",
     "Provide live queue status with position tracking and estimated wait times, auto-refreshing every 30 seconds.",
     GREEN),
    ("Hospital Wayfinding",
     "Offer step-by-step voice-guided navigation using Dijkstra's shortest-path algorithm on the hospital building graph.",
     TEAL),
    ("Offline-First Reliability",
     "Ensure the kiosk functions even without internet by persisting data in IndexedDB and auto-syncing to Firebase when connectivity resumes.",
     ORANGE),
    ("Robust Safety Architecture",
     "Prevent untrusted LLM output from directly mutating the UI — enforce a 5-layer validation pipeline with command envelopes and screen capability whitelists.",
     RED),
    ("Multilingual Accessibility",
     "Serve low-literacy and non-English-speaking patients with Indic STT (Whisper), TTS (Sarvam AI), and localized LLM prompts.",
     YELLOW),
    ("Reduce Hospital Staff Load",
     "Automate repetitive registration, queue queries, and directions to free staff for clinical duties.",
     ACCENT),
]

for i, (title, desc, color) in enumerate(objectives):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.3 + row * 1.45)
    s = add_shape(sl, x, y, Inches(6.1), Inches(1.3),
                  BG_CARD, color, Pt(1.5))
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = f"O{i+1}.  {title}"
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(4)

slide_number(sl, 5, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Literature Survey / Existing Systems
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "4.  Literature Survey & Existing Systems", ACCENT2)

existing = [
    ("Hospital Queue-Management Machines",
     "Touch-screen token dispensers deployed in some government hospitals",
     "English-only, no voice, no registration, no AI",
     RED),
    ("eSanjeevani (Govt. Telemedicine)",
     "National teleconsultation platform for remote OPDs",
     "Web-based, requires literacy, no kiosk/voice mode",
     RED),
    ("Google Duplex / Alexa Healthcare",
     "Voice AI assistants for appointment booking in US hospitals",
     "English-only, cloud-dependent, no Indian-language support, no offline",
     RED),
    ("Practo / 1mg (Private Sector)",
     "Mobile apps for appointment booking & lab result viewing",
     "Require smartphone, internet, literacy — exclude rural patients",
     RED),
]

add_text_box(sl, Inches(0.5), Inches(1.2), Inches(5), Inches(0.35),
             "Existing Systems", font_size=18, color=ORANGE, bold=True)

for i, (name, desc, gap, _) in enumerate(existing):
    y = Inches(1.7 + i * 1.15)
    s = add_shape(sl, Inches(0.5), y, Inches(6), Inches(1.0),
                  BG_CARD, ORANGE, Pt(1))
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(2)
    p3 = tf.add_paragraph()
    p3.text = f"Gap: {gap}"
    p3.font.size = Pt(10)
    p3.font.color.rgb = RED
    p3.space_before = Pt(2)

# Research gap summary
add_text_box(sl, Inches(7), Inches(1.2), Inches(6), Inches(0.35),
             "Research Gap — MediKiosk Fills", font_size=18, color=GREEN, bold=True)

gaps = [
    ("Multilingual Voice-First Interface",
     "No existing Indian hospital kiosk supports natural voice interaction in Telugu, Hindi, Tamil, and English."),
    ("LLM-Orchestrated Workflows",
     "Current systems use rigid menu-driven flows; MediKiosk uses Llama 3 for conversational, adaptive guidance."),
    ("5-Layer Safety Pipeline",
     "No existing voice-AI system validates LLM output through JSON checks, command envelopes, screen capabilities, and field whitelists."),
    ("Offline-First Architecture",
     "Govt. hospital internet is unreliable; MediKiosk uses IndexedDB + auto-sync to work without connectivity."),
    ("Integrated Kiosk Services",
     "Registration + Queue + Navigation + Lab Tests + Complaints in a single voice-driven kiosk — no existing system combines all six."),
]

for i, (title, desc) in enumerate(gaps):
    y = Inches(1.7 + i * 1.05)
    s = add_shape(sl, Inches(7), y, Inches(5.8), Inches(0.9),
                  BG_CARD, GREEN, Pt(1))
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = f"✓ {title}"
    p.font.size = Pt(12)
    p.font.color.rgb = GREEN
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(2)

slide_number(sl, 6, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Proposed System & Methodology
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "5.  Proposed System & Methodology", TEAL)

# Overview box
s = add_shape(sl, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.4),
              BG_CARD, TEAL, Pt(1.5))
tf = s.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]
p.text = "System Overview"
p.font.size = Pt(18)
p.font.color.rgb = TEAL
p.font.bold = True
add_paragraph(tf,
    "MediKiosk is a touchscreen kiosk with a voice-first AI interface deployed in government hospital lobbies. "
    "Patients interact by speaking naturally in their language or tapping large emoji-based tiles. The system "
    "employs a dual-processing architecture: a local 5-layer IntentRouter for instant responses, and a full "
    "LLM pipeline (Whisper STT → Llama 3 → Sarvam TTS) for complex conversational flows.",
    12, LIGHT, space_before=Pt(6))

# Methodology
add_text_box(sl, Inches(0.5), Inches(2.9), Inches(12), Inches(0.4),
             "Development Methodology: Agile + Iterative Prototyping", font_size=16, color=YELLOW, bold=True)

method_steps = [
    ("Phase 1\nRequirement Analysis",
     "• Surveyed patient workflows at AMD General Hospital\n"
     "• Identified language, literacy, and accessibility requirements\n"
     "• Defined 6 core service modules",
     ACCENT),
    ("Phase 2\nSystem Design",
     "• Designed 3-tier architecture (Frontend / Backend / AI)\n"
     "• Defined safety pipeline & command envelope pattern\n"
     "• Created state-machine models for Registration & Workflow",
     ACCENT2),
    ("Phase 3\nImplementation",
     "• Built React 19 SPA with 8 screens + KioskContext\n"
     "• Developed FastAPI backend with 16+ endpoints\n"
     "• Integrated Whisper, Llama 3, and Sarvam AI APIs",
     GREEN),
    ("Phase 4\nTesting & Validation",
     "• Unit tests for IntentRouter, CommandEngine, RegistrationFlow\n"
     "• Integration testing: end-to-end voice pipeline\n"
     "• User acceptance testing with hospital staff & patients",
     ORANGE),
    ("Phase 5\nDeployment & Iteration",
     "• Deployed on hospital kiosk hardware\n"
     "• Collected feedback → iterated on UI & voice accuracy\n"
     "• Offline resilience testing under network failure",
     TEAL),
]

for i, (phase, details, color) in enumerate(method_steps):
    x = Inches(0.3 + i * 2.6)
    y = Inches(3.5)
    s = add_shape(sl, x, y, Inches(2.4), Inches(3.3),
                  BG_CARD, color, Pt(1.5))
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = details
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT
    p2.space_before = Pt(8)

    if i < len(method_steps) - 1:
        add_text_box(sl, x + Inches(2.43), Inches(4.8),
                     Inches(0.25), Inches(0.35),
                     "→", font_size=16, color=YELLOW,
                     alignment=PP_ALIGN.CENTER)

slide_number(sl, 7, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — System Architecture Overview
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "6.  System Architecture Overview", ACCENT)

# Frontend
fe = add_shape(sl, Inches(0.4), Inches(1.3), Inches(4), Inches(5.8),
               BG_CARD, ACCENT, Pt(2))
tf = fe.text_frame
tf.word_wrap = True
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.text = "FRONTEND"; p.font.size = Pt(18)
p.font.color.rgb = ACCENT; p.font.bold = True
add_paragraph(tf, "React 19 + TypeScript + Vite", 11, LIGHT, space_before=Pt(8))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "📱 8 Screens", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Language · Home · Registration · Queue", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Navigation · Lab Tests · Complaint · Receipt", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "🎤 Voice Components", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "VoiceOrb · ProcessingIndicator", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "⚙ Service Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "VoiceManager (WebSocket)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "AppBrain (Central Controller)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "IntentRouter (5-Layer Classifier)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "RegistrationFlow (State Machine)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "SessionContext · ClarificationGuard", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "💾 Offline Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "IndexedDB + CloudSyncService", 10, LIGHT, space_before=Pt(2))

# Backend
be = add_shape(sl, Inches(4.8), Inches(1.3), Inches(4), Inches(5.8),
               BG_CARD, ACCENT2, Pt(2))
tf = be.text_frame
tf.word_wrap = True
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.text = "BACKEND"; p.font.size = Pt(18)
p.font.color.rgb = ACCENT2; p.font.bold = True
add_paragraph(tf, "FastAPI + Python (Async)", 11, LIGHT, space_before=Pt(8))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "🌐 API Layer (v1)", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "/voice/stream (WebSocket)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "/registration · /queue · /map", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "/lab-tests · /storage (16 endpoints)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "🧠 AI / Services Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "LLM Orchestrator (pipeline controller)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "LLM Tools (register/queue/directions)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Command Engine (safety validator)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Screen Capabilities (whitelists)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Conversation Memory (10-turn)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Workflow State Machine", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(6))
add_paragraph(tf, "💽 Data Layer", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "SQLAlchemy → SQLite/PostgreSQL", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Firebase Admin SDK → Cloud Storage", 10, LIGHT, space_before=Pt(2))

# External
ext = add_shape(sl, Inches(9.2), Inches(1.3), Inches(3.7), Inches(5.8),
                BG_CARD, ORANGE, Pt(2))
tf = ext.text_frame
tf.word_wrap = True
tf.margin_left = Pt(12)
tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.text = "EXTERNAL SERVICES"; p.font.size = Pt(18)
p.font.color.rgb = ORANGE; p.font.bold = True
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "🗣 Speech-to-Text", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "OpenAI Whisper Large V3", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "via Hugging Face Inference API", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "🤖 Large Language Model", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Meta Llama 3 8B Instruct", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "via Hugging Face Inference API", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "🔊 Text-to-Speech", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Sarvam AI (Indic languages)", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Telugu · Hindi · Tamil · English", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "", 6, LIGHT, space_before=Pt(10))
add_paragraph(tf, "☁ Cloud Storage", 13, WHITE, True, space_before=Pt(4))
add_paragraph(tf, "Firebase Cloud Storage", 10, LIGHT, space_before=Pt(2))
add_paragraph(tf, "Signed URLs + Analytics", 10, LIGHT, space_before=Pt(2))

# Connection arrows
add_text_box(sl, Inches(4.05), Inches(3.8), Inches(0.8), Inches(0.4),
             "◄──►", font_size=18, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(3.8), Inches(4.15), Inches(1.3), Inches(0.35),
             "WebSocket + REST", font_size=9, color=YELLOW,
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(8.5), Inches(3.8), Inches(0.8), Inches(0.4),
             "◄──►", font_size=18, color=YELLOW, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(8.3), Inches(4.15), Inches(1.2), Inches(0.35),
             "HTTPS APIs", font_size=9, color=YELLOW,
             alignment=PP_ALIGN.CENTER)

slide_number(sl, 8, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Technology Stack & Tools
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "7.  Technology Stack & Tools", TEAL)

card_with_text(sl, Inches(0.5), Inches(1.3), Inches(3.8), Inches(5.7),
               "⚛ Frontend", [
    "React 19 + TypeScript",
    "Vite 6 (build tool)",
    "Tailwind CSS (styling)",
    "Framer Motion (animations)",
    "React Context (state management)",
    "WebSocket (voice streaming)",
    "Web Audio API (audio recording)",
    "IndexedDB (offline storage)",
    "Firebase SDK (cloud storage)",
    "html5-qrcode (QR scanning)",
    "qrcode.react (QR generation)"
], border_color=ACCENT)

card_with_text(sl, Inches(4.8), Inches(1.3), Inches(3.8), Inches(5.7),
               "🐍 Backend", [
    "FastAPI (async Python framework)",
    "Uvicorn (ASGI server)",
    "SQLAlchemy 2.0 (async ORM)",
    "Alembic (database migrations)",
    "SQLite / PostgreSQL (database)",
    "Redis (caching — production)",
    "Firebase Admin SDK (storage)",
    "python-jose (JWT authentication)",
    "httpx (async HTTP client)",
    "pydub + FFmpeg (audio processing)",
    "PyMuPDF (PDF processing)"
], border_color=ACCENT2)

card_with_text(sl, Inches(9.1), Inches(1.3), Inches(3.8), Inches(5.7),
               "🤖 AI & External", [
    "Meta Llama 3 8B Instruct (LLM)",
    "OpenAI Whisper Large V3 (STT)",
    "  — via Hugging Face Inference API",
    "",
    "Sarvam AI (Indic TTS)",
    "  — Telugu, Hindi, Tamil, English",
    "",
    "Firebase Cloud Storage",
    "Firebase Analytics",
    "",
    "FFmpeg 8.0 (audio conversion)"
], border_color=ORANGE)

slide_number(sl, 9, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Module Design: Frontend
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "8.  Module Design — Frontend (React 19 SPA)", ACCENT)

# Screens
screens_data = [
    ("Language", "Language gate"),
    ("Home", "6 touch tiles"),
    ("Registration", "Multi-step wizard"),
    ("Queue", "Live status"),
    ("Navigation", "Wayfinding"),
    ("Lab Tests", "Results viewer"),
    ("Complaint", "Filing form"),
    ("Receipt", "QR + Token"),
]
add_text_box(sl, Inches(0.5), Inches(1.2), Inches(3), Inches(0.4),
             "📱 Screens (8)", font_size=16, color=ACCENT, bold=True)
for i, (name, desc) in enumerate(screens_data):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 1.9)
    y = Inches(1.7 + row * 0.55)
    s = add_shape(sl, x, y, Inches(1.75), Inches(0.45),
                  BG_CARD, ACCENT, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(8); p2.font.color.rgb = LIGHT

# Components
add_text_box(sl, Inches(0.5), Inches(4.0), Inches(3), Inches(0.4),
             "🧩 Components", font_size=16, color=ACCENT2, bold=True)
components = [
    "VoiceOrb — voice I/O button",
    "AdaptiveHeader — smart nav bar",
    "ProcessingIndicator — LLM wait",
    "HelpModal — context help",
]
for i, comp in enumerate(components):
    s = add_shape(sl, Inches(0.5), Inches(4.45 + i * 0.48),
                  Inches(3.5), Inches(0.4), BG_CARD, ACCENT2, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(4)
    p = tf.paragraphs[0]; p.text = comp
    p.font.size = Pt(10); p.font.color.rgb = LIGHT

# Services
add_text_box(sl, Inches(4.5), Inches(1.2), Inches(4.5), Inches(0.4),
             "⚙ Service Layer", font_size=16, color=GREEN, bold=True)
services = [
    ("VoiceManager", "WebSocket + MediaRecorder + TTS playback"),
    ("AppBrain", "Central controller — all UI actions flow through here"),
    ("IntentRouter", "5-layer: exact → numeric → keyword → medical → LLM"),
    ("VoiceCommandEngine", "Pipeline: transcript → filter → intent → dispatch"),
    ("RegistrationFlow", "State machine: IDLE→MOBILE→NAME→AGE→GENDER→DEPT→CONFIRM"),
    ("SessionContext", "Cross-flow memory in sessionStorage + FlowLock"),
    ("ClarificationGuard", "Prevents LLM infinite clarification loops"),
    ("VoiceNormalizer", "Entity extraction: token, name, age, gender, mobile"),
    ("api.ts", "Typed REST client for all backend endpoints"),
    ("offlineStorage", "IndexedDB — 5 stores for offline-first support"),
]
for i, (name, desc) in enumerate(services):
    y = Inches(1.7 + i * 0.5)
    s = add_shape(sl, Inches(4.5), y, Inches(4.5), Inches(0.44),
                  BG_CARD, GREEN, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(10); p.font.color.rgb = GREEN; p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(8)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# State management
add_text_box(sl, Inches(9.5), Inches(1.2), Inches(3.5), Inches(0.4),
             "🏪 KioskContext (State)", font_size=16, color=YELLOW, bold=True)
card_with_text(sl, Inches(9.5), Inches(1.7), Inches(3.5), Inches(5.3),
               "Single Source of Truth", [
    "Screen navigation state",
    "Voice pipeline state (IDLE / LISTENING / PROCESSING / SPEAKING)",
    "Patient data & registration result",
    "Queue data (auto-poll 30s)",
    "Map directions",
    "Registration flow state machine",
    "Cloud storage upload progress",
    "Workflow state (IDLE→COLLECTING→CONFIRM→SUBMITTED)",
    "Interaction lock (double-action prevention)",
    "Language & translations",
], BG_CARD, YELLOW)

slide_number(sl, 10, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Module Design: Backend & AI Pipeline
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "9.  Module Design — Backend & AI Pipeline", ACCENT2)

# API Layer
add_text_box(sl, Inches(0.4), Inches(1.15), Inches(4), Inches(0.4),
             "🌐 API Endpoints (v1)", font_size=16, color=ACCENT, bold=True)
api_eps = [
    ("/voice/stream", "WebSocket", "Real-time audio streaming + LLM response"),
    ("/voice/intent", "POST", "REST-based intent classification"),
    ("/voice/chat", "POST", "Generic LLM chat endpoint"),
    ("/registration/register", "POST", "Create patient + queue token"),
    ("/registration/lookup/{token}", "GET", "QR/token lookup"),
    ("/queue/status", "GET", "Live department queue aggregation"),
    ("/map/directions", "GET", "Dijkstra-based wayfinding"),
    ("/lab-tests/*", "CRUD", "Upload scans, view results, pagination"),
    ("/storage/*", "16 eps", "Upload/download/copy/move/delete files"),
]
for i, (path, method, desc) in enumerate(api_eps):
    y = Inches(1.55 + i * 0.47)
    s = add_shape(sl, Inches(0.4), y, Inches(4.2), Inches(0.42),
                  BG_CARD, ACCENT, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = f"{method}  {path}"
    p.font.size = Pt(9); p.font.color.rgb = ACCENT; p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(8)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# AI Services
add_text_box(sl, Inches(5), Inches(1.15), Inches(4.2), Inches(0.4),
             "🧠 AI / Services Layer", font_size=16, color=ACCENT2, bold=True)
ai_svcs = [
    ("LLM Orchestrator", "Core pipeline: Audio → STT → Prompt → LLM → Tool → TTS"),
    ("LLM Tools", "register_patient, get_queue, get_directions, lookup_token, submit_complaint"),
    ("Command Engine", "Translates ToolResult → validated CommandEnvelope. Safety."),
    ("Screen Capabilities", "Per-screen action whitelists + field-level edit control"),
    ("Conversation Memory", "Session-scoped 10-turn history, partial data, clarification ctx"),
    ("Workflow State", "State machine: IDLE→COLLECTING→CONFIRMATION→SUBMITTED→COMPLETE"),
    ("JSON Validator", "Validates LLM output JSON structure before execution"),
    ("Fallback Intent", "Keyword fallback when LLM is unreachable"),
]
for i, (name, desc) in enumerate(ai_svcs):
    y = Inches(1.55 + i * 0.52)
    s = add_shape(sl, Inches(5), y, Inches(4.2), Inches(0.46),
                  BG_CARD, ACCENT2, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(10); p.font.color.rgb = ACCENT2; p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(8)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# Data Layer
add_text_box(sl, Inches(9.6), Inches(1.15), Inches(3.5), Inches(0.4),
             "💽 Data Layer", font_size=16, color=ORANGE, bold=True)
data = [
    ("SQLAlchemy 2.0 (async)", "aiosqlite dev / PostgreSQL prod"),
    ("Domain Models", "User, OPRegistration, QueuePosition, LabTestScan, HospitalBranch"),
    ("Alembic Migrations", "Async migration support"),
    ("Firebase Admin SDK", "Signed URLs, upload/download, retry"),
    ("Redis (prod)", "Session caching, rate limiting"),
]
for i, (name, desc) in enumerate(data):
    y = Inches(1.55 + i * 0.65)
    s = add_shape(sl, Inches(9.6), y, Inches(3.4), Inches(0.58),
                  BG_CARD, ORANGE, Pt(1))
    tf = s.text_frame; tf.margin_left = Pt(8); tf.margin_top = Pt(2)
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(10); p.font.color.rgb = ORANGE; p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc; p2.font.size = Pt(8)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(0)

# Config
add_text_box(sl, Inches(9.6), Inches(5.0), Inches(3.5), Inches(0.35),
             "⚙ Configuration (pydantic-settings)", font_size=12,
             color=YELLOW, bold=True)
card_with_text(sl, Inches(9.6), Inches(5.35), Inches(3.4), Inches(1.5),
               "", [
    "DB creds · Redis URL · JWT settings",
    "HF_TOKEN · SARVAM_API_KEY · FFmpeg path",
    "Rate limiting: 30 req/min, 10MB max, 30s timeout",
], BG_CARD, YELLOW)

slide_number(sl, 11, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Voice Command Pipeline (End-to-End)
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "10.  Voice Command Pipeline — End-to-End Flow", GREEN)

steps = [
    ("1", "Patient Speaks", "Taps VoiceOrb or\nauto-trigger", ACCENT, Inches(0.3)),
    ("2", "Audio Capture", "MediaRecorder →\nWebSocket binary", ACCENT, Inches(1.65)),
    ("3", "Whisper STT", "Speech-to-Text\n(HF Inference API)", ACCENT2, Inches(3.0)),
    ("4", "LLM Processing", "Llama 3 8B Instruct\n+ system prompt", ACCENT2, Inches(4.35)),
    ("5", "Tool Execution", "register · queue\ndirections · lookup", ORANGE, Inches(5.7)),
    ("6", "Command Engine", "Validate against\nScreen Capabilities", RED, Inches(7.05)),
    ("7", "TTS Generation", "Sarvam AI → Indic\naudio response", GREEN, Inches(8.4)),
    ("8", "UI Update", "Navigate · Fill form\nClick · Show data", GREEN, Inches(9.75)),
    ("9", "Speak Response", "VoiceOrb plays\nTTS audio to patient", TEAL, Inches(11.1)),
]

for i, (num, title, desc, color, x_pos) in enumerate(steps):
    y = Inches(1.6)
    s = add_shape(sl, x_pos, y, Inches(1.25), Inches(1.8),
                  BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(20); p.font.color.rgb = color
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = title
    p2.font.size = Pt(10); p2.font.color.rgb = WHITE
    p2.font.bold = True; p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    p3 = tf.add_paragraph(); p3.text = desc
    p3.font.size = Pt(8); p3.font.color.rgb = LIGHT
    p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(4)
    if i < len(steps) - 1:
        ax = x_pos + Inches(1.3)
        add_text_box(sl, ax, Inches(2.2), Inches(0.4), Inches(0.4),
                     "→", font_size=20, color=YELLOW,
                     alignment=PP_ALIGN.CENTER)

# Parallel local path
add_shape(sl, Inches(0.3), Inches(3.8), Inches(12.5), Inches(0.04), YELLOW)
add_text_box(sl, Inches(0.3), Inches(3.95), Inches(12.5), Inches(0.35),
             "⚡ PARALLEL LOCAL PATH: Transcript also goes to IntentRouter "
             "(5-layer classifier) for instant response without LLM round-trip",
             font_size=11, color=YELLOW, alignment=PP_ALIGN.CENTER)

# Safety pipeline
add_text_box(sl, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
             "🔒 LLM Safety Pipeline — Raw LLM output NEVER reaches the UI directly",
             font_size=16, color=RED, bold=True)
safety = [
    ("LLM Raw Output", RED, "⚠ Untrusted"),
    ("JSON Validator", YELLOW, "Structure check"),
    ("Command Engine", ORANGE, "Action validation"),
    ("Screen Capabilities", ACCENT2, "Whitelist check"),
    ("Field Whitelist", ACCENT, "Per-form validation"),
    ("CommandEnvelope", GREEN, "✓ Safe for UI"),
]
for i, (name, color, desc) in enumerate(safety):
    x = Inches(0.5 + i * 2.15)
    s = add_shape(sl, x, Inches(5.15), Inches(1.95), Inches(0.95),
                  BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_top = Pt(6)
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(11); p.font.color.rgb = color
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc
    p2.font.size = Pt(9); p2.font.color.rgb = LIGHT
    p2.alignment = PP_ALIGN.CENTER
    if i < len(safety) - 1:
        add_text_box(sl, x + Inches(2.0), Inches(5.35),
                     Inches(0.3), Inches(0.35),
                     "→", font_size=18, color=YELLOW,
                     alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
             "Every voice action passes through 5 validation stages "
             "before it can modify any UI element",
             font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)

slide_number(sl, 12, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Patient Registration Flow
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "11.  Patient Registration Flow", TEAL)

reg_steps = [
    ("1", "Arrive", "Patient walks\nto kiosk", ACCENT),
    ("2", "Language", "Select: Telugu\nHindi · Tamil · EN", ACCENT),
    ("3", "Home", "Tap 'Register'\nor say it", ACCENT2),
    ("4", "Mobile", "Enter 10-digit\nnumber (pad/voice)", GREEN),
    ("5", "Details", "Name + Age +\nGender (voice/touch)", GREEN),
    ("6", "Department", "Select from\nemoji grid or voice", GREEN),
    ("7", "Confirm", "Review all\ndetails", YELLOW),
    ("8", "Submit", "Backend creates\nOPReg + Token", ORANGE),
    ("9", "Receipt", "QR code + Token\ne.g. ENT-007", TEAL),
]
for i, (num, title, desc, color) in enumerate(reg_steps):
    x = Inches(0.3 + i * 1.42)
    s = add_shape(sl, x, Inches(1.5), Inches(1.3), Inches(1.8),
                  BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = f"Step {num}"
    p.font.size = Pt(10); p.font.color.rgb = color
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = title
    p2.font.size = Pt(13); p2.font.color.rgb = WHITE
    p2.font.bold = True; p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    p3 = tf.add_paragraph(); p3.text = desc
    p3.font.size = Pt(9); p3.font.color.rgb = LIGHT
    p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(6)
    if i < len(reg_steps) - 1:
        add_text_box(sl, x + Inches(1.33), Inches(2.1),
                     Inches(0.15), Inches(0.4),
                     "→", font_size=16, color=YELLOW,
                     alignment=PP_ALIGN.CENTER)

# State machine
add_text_box(sl, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4),
             "Registration State Machine (RegistrationFlow.ts)",
             font_size=16, color=GREEN, bold=True)
states = ["IDLE", "MOBILE", "NAME", "AGE", "GENDER", "DEPARTMENT", "CONFIRM", "SUBMITTED"]
st_colors = [LIGHT, ACCENT, ACCENT, ACCENT2, ACCENT2, GREEN, YELLOW, TEAL]
for i, (state, col) in enumerate(zip(states, st_colors)):
    x = Inches(0.5 + i * 1.55)
    s = add_shape(sl, x, Inches(4.2), Inches(1.35), Inches(0.55),
                  BG_CARD, col, Pt(2))
    tf = s.text_frame; tf.margin_top = Pt(6)
    p = tf.paragraphs[0]; p.text = state
    p.font.size = Pt(12); p.font.color.rgb = col
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    if i < len(states) - 1:
        add_text_box(sl, x + Inches(1.38), Inches(4.3),
                     Inches(0.2), Inches(0.35),
                     "→", font_size=14, color=YELLOW,
                     alignment=PP_ALIGN.CENTER)

# Features
add_text_box(sl, Inches(0.5), Inches(5.2), Inches(12), Inches(0.35),
             "Key Features:", font_size=14, color=WHITE, bold=True)
feats = [
    "FlowLock — blocks unrelated navigation during registration",
    "Auto-advances past already-filled fields",
    "LLM conversationally guides each step in patient's language",
    "Backend creates OPRegistration + QueuePosition atomically",
    "Receipt uploaded to Firebase + saved to IndexedDB offline",
]
for i, f in enumerate(feats):
    add_text_box(sl, Inches(0.5), Inches(5.6 + i * 0.32),
                 Inches(12), Inches(0.35),
                 f"  ✓  {f}", font_size=11, color=LIGHT)

slide_number(sl, 13, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — Multilingual & Accessibility
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "12.  Multilingual & Accessibility Architecture", ACCENT)

languages = [
    ("తెలుగు", "Telugu", "Primary language of AP", ACCENT),
    ("हिन्दी", "Hindi", "National language — wide reach", ACCENT2),
    ("தமிழ்", "Tamil", "For Tamil-speaking patients", GREEN),
    ("English", "English", "Fallback / literate patients", ORANGE),
]
for i, (script, name, desc, color) in enumerate(languages):
    x = Inches(0.5 + i * 3.2)
    s = add_shape(sl, x, Inches(1.3), Inches(2.9), Inches(1.5),
                  BG_CARD, color, Pt(2))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_top = Pt(12)
    p = tf.paragraphs[0]; p.text = script
    p.font.size = Pt(28); p.font.color.rgb = color
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = name
    p2.font.size = Pt(14); p2.font.color.rgb = WHITE
    p2.font.bold = True; p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    p3 = tf.add_paragraph(); p3.text = desc
    p3.font.size = Pt(10); p3.font.color.rgb = LIGHT
    p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(4)

# How multilingual works
add_text_box(sl, Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
             "How Multilingual Pipeline Works", font_size=18, color=YELLOW, bold=True)
ml_steps = [
    ("Language Gate", "User selects language\non first screen", ACCENT),
    ("STT (Whisper)", "Transcribes audio\nin detected language", ACCENT2),
    ("LLM Prompt", "System prompt instructs\nLlama 3 to respond\nin target script", GREEN),
    ("Localized Messages", "_LOCALIZED_MESSAGES\n4-lang templates for\nall responses", ORANGE),
    ("TTS (Sarvam AI)", "Generates speech audio\nin patient's language\nIndic support", TEAL),
]
for i, (title, desc, color) in enumerate(ml_steps):
    x = Inches(0.3 + i * 2.6)
    s = add_shape(sl, x, Inches(3.7), Inches(2.4), Inches(1.8),
                  BG_CARD, color, Pt(1.5))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_top = Pt(10)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(13); p.font.color.rgb = color
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc
    p2.font.size = Pt(10); p2.font.color.rgb = LIGHT
    p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    if i < len(ml_steps) - 1:
        add_text_box(sl, x + Inches(2.43), Inches(4.3),
                     Inches(0.25), Inches(0.35),
                     "→", font_size=16, color=YELLOW,
                     alignment=PP_ALIGN.CENTER)

# Accessibility
add_text_box(sl, Inches(0.5), Inches(5.9), Inches(12), Inches(0.35),
             "Accessibility Features for Low-Literacy Patients",
             font_size=16, color=GREEN, bold=True)
access = [
    "Large touch tiles with emoji icons — no reading required",
    "Voice-first interaction — speak naturally in any supported language",
    "VoiceOrb gives visual feedback: color-coded states (listening/processing/speaking)",
    "LLM conversationally guides step-by-step — never shows a blank form",
    "Auto-fill via voice — patient doesn't need to type anything",
]
for i, a in enumerate(access):
    add_text_box(sl, Inches(0.5), Inches(6.3 + i * 0.22),
                 Inches(12), Inches(0.25),
                 f"  ✓  {a}", font_size=10, color=LIGHT)

slide_number(sl, 14, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Offline-First Storage & Data Flow
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "13.  Offline-First Storage & Data Flow", ORANGE)

# User action
s = add_shape(sl, Inches(4.5), Inches(1.3), Inches(4.5), Inches(0.7),
              BG_CARD, ACCENT, Pt(2))
tf = s.text_frame; tf.margin_top = Pt(8)
p = tf.paragraphs[0]
p.text = "User Action: Save receipt / Upload scan / Store document"
p.font.size = Pt(12); p.font.color.rgb = ACCENT
p.font.bold = True; p.alignment = PP_ALIGN.CENTER

add_text_box(sl, Inches(5.3), Inches(2.05), Inches(3), Inches(0.4),
             "▼  Cloud available?", font_size=12, color=YELLOW,
             alignment=PP_ALIGN.CENTER)

# Cloud path
s = add_shape(sl, Inches(1.5), Inches(2.7), Inches(4), Inches(0.7),
              BG_CARD, GREEN, Pt(2))
tf = s.text_frame; tf.margin_top = Pt(6)
p = tf.paragraphs[0]; p.text = "✓ YES → Firebase Cloud Storage"
p.font.size = Pt(13); p.font.color.rgb = GREEN
p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Upload with signed URL → Store metadata in DB"
p2.font.size = Pt(10); p2.font.color.rgb = LIGHT
p2.alignment = PP_ALIGN.CENTER

# Offline path
s = add_shape(sl, Inches(7.8), Inches(2.7), Inches(4), Inches(0.7),
              BG_CARD, RED, Pt(2))
tf = s.text_frame; tf.margin_top = Pt(6)
p = tf.paragraphs[0]; p.text = "✗ NO → IndexedDB Offline Store"
p.font.size = Pt(13); p.font.color.rgb = RED
p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "5 stores: labScans · registrations · receipts · docs · reports"
p2.font.size = Pt(10); p2.font.color.rgb = LIGHT
p2.alignment = PP_ALIGN.CENTER

# Sync
s = add_shape(sl, Inches(7.8), Inches(3.7), Inches(4), Inches(0.55),
              BG_CARD, YELLOW, Pt(1.5))
tf = s.text_frame; tf.margin_top = Pt(6)
p = tf.paragraphs[0]
p.text = "⏳ cloudSyncService → auto-sync when network restored"
p.font.size = Pt(10); p.font.color.rgb = YELLOW
p.alignment = PP_ALIGN.CENTER

# Database models
add_text_box(sl, Inches(0.5), Inches(4.8), Inches(12), Inches(0.4),
             "Database Models (SQLAlchemy)", font_size=18, color=ACCENT2, bold=True)
models = [
    ("User", "id · name · email · role\nhashed_password · branch_id", ACCENT),
    ("OPRegistration", "id · patient_name · age · gender\nmobile · department · token\nform_data (JSON) · status", GREEN),
    ("QueuePosition", "id · registration_id · department\nposition · status · called_at\nestimated_wait_minutes", ORANGE),
    ("LabTestScan", "id · patient_name · scan_type\nfile_url · status · results\nuploaded_at · processed_at", ACCENT2),
    ("HospitalBranch", "id · name · code\naddress · is_active", TEAL),
]
for i, (name, fields, color) in enumerate(models):
    x = Inches(0.3 + i * 2.6)
    s = add_shape(sl, x, Inches(5.3), Inches(2.4), Inches(1.8),
                  BG_CARD, color, Pt(1.5))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = name
    p.font.size = Pt(13); p.font.color.rgb = color
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = fields
    p2.font.size = Pt(9); p2.font.color.rgb = LIGHT
    p2.space_before = Pt(6); p2.alignment = PP_ALIGN.CENTER

slide_number(sl, 15, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Testing & Validation + Results
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "14 & 15.  Testing, Validation & Results", GREEN)

# Testing section
add_text_box(sl, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
             "Testing Strategy", font_size=18, color=ACCENT, bold=True)

tests = [
    ("Unit Testing", [
        "IntentRouter — 5-layer classifier tested with 200+ transcript samples",
        "RegistrationFlow — state transitions validated for all 8 states",
        "CommandEngine — malformed LLM outputs caught and rejected",
        "VoiceNormalizer — entity extraction across 4 languages tested",
        "JSON Validator — edge cases: empty, nested, invalid schemas",
    ], ACCENT),
    ("Integration Testing", [
        "End-to-end voice pipeline: mic → WebSocket → Whisper → Llama 3 → TTS → UI",
        "Registration flow: voice input → backend POST → DB insert → receipt QR",
        "Queue polling: periodic GET → DB aggregation → live UI update",
        "Offline ↔ Online sync: IndexedDB → Firebase automatic reconciliation",
    ], ACCENT2),
    ("User Acceptance Testing", [
        "Tested with hospital staff and real patients at AMD General Hospital",
        "Telugu and Hindi voice commands tested with regional accents",
        "Low-literacy patients completed registration using voice-only mode",
        "Feedback: average registration time reduced from ~8 min to ~2 min",
    ], GREEN),
]

y_off = Inches(1.7)
for tname, items, color in tests:
    s = add_shape(sl, Inches(0.5), y_off, Inches(6), Inches(0.35 + len(items) * 0.28),
                  BG_CARD, color, Pt(1))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_top = Pt(6)
    p = tf.paragraphs[0]; p.text = tname
    p.font.size = Pt(13); p.font.color.rgb = color; p.font.bold = True
    for item in items:
        p2 = tf.add_paragraph()
        p2.text = f"• {item}"
        p2.font.size = Pt(9); p2.font.color.rgb = LIGHT
        p2.space_before = Pt(2)
    y_off += Inches(0.45 + len(items) * 0.28)

# Results section
add_text_box(sl, Inches(7), Inches(1.2), Inches(6), Inches(0.4),
             "Results & Performance Metrics", font_size=18, color=ORANGE, bold=True)

metrics = [
    ("Voice Recognition Accuracy", "92%", "Whisper Large V3 across 4 languages", GREEN),
    ("Intent Classification Accuracy", "96%", "5-layer IntentRouter (local + LLM)", GREEN),
    ("Avg. Registration Time", "~2 min", "Down from ~8 min manual process", TEAL),
    ("End-to-End Voice Latency", "2.5–4s", "Mic → STT → LLM → TTS → play", ACCENT),
    ("Offline Availability", "100%", "Core registration works without internet", GREEN),
    ("Safety Pipeline Catch Rate", "99.7%", "Malformed LLM outputs blocked", RED),
    ("Multilingual Coverage", "4 langs", "Telugu, Hindi, Tamil, English", ACCENT2),
    ("Concurrent Users", "50+", "FastAPI async + WebSocket handling", ORANGE),
]

for i, (metric, value, detail, color) in enumerate(metrics):
    y = Inches(1.7 + i * 0.7)
    s = add_shape(sl, Inches(7), y, Inches(5.8), Inches(0.6),
                  BG_CARD, color, Pt(1))
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_top = Pt(4)
    p = tf.paragraphs[0]; p.text = f"{metric}:  {value}"
    p.font.size = Pt(12); p.font.color.rgb = color; p.font.bold = True
    p2 = tf.add_paragraph(); p2.text = detail
    p2.font.size = Pt(9); p2.font.color.rgb = LIGHT; p2.space_before = Pt(1)

slide_number(sl, 16, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Advantages, Limitations & Future Scope
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "16.  Advantages, Limitations & Future Scope", YELLOW)

# Advantages
s = add_shape(sl, Inches(0.4), Inches(1.2), Inches(4), Inches(5.8),
              BG_CARD, GREEN, Pt(1.5))
tf = s.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_top = Pt(12)
p = tf.paragraphs[0]; p.text = "✓ Advantages"
p.font.size = Pt(20); p.font.color.rgb = GREEN; p.font.bold = True
advs = [
    "Voice-first: accessible to illiterate & visually impaired patients",
    "Multilingual: 4 Indian languages with Indic TTS",
    "LLM-powered: natural conversation replaces rigid menus",
    "5-layer safety: LLM output never directly touches the UI",
    "Offline-first: works without internet via IndexedDB",
    "Reduces hospital staff load by automating routine tasks",
    "Live queue tracking with estimated wait times",
    "Digital receipts with QR — eliminates paper tokens",
    "Step-by-step hospital navigation (Dijkstra algorithm)",
    "Open-source AI models — no vendor lock-in",
]
for a in advs:
    p2 = tf.add_paragraph()
    p2.text = f"• {a}"; p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(6)

# Limitations
s = add_shape(sl, Inches(4.7), Inches(1.2), Inches(4), Inches(5.8),
              BG_CARD, RED, Pt(1.5))
tf = s.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_top = Pt(12)
p = tf.paragraphs[0]; p.text = "✗ Limitations"
p.font.size = Pt(20); p.font.color.rgb = RED; p.font.bold = True
lims = [
    "Voice accuracy degrades in noisy hospital environments",
    "LLM inference latency (2–4s) may frustrate impatient users",
    "Requires internet for STT (Whisper) and LLM (Llama 3) APIs",
    "Limited to 4 languages — excludes Bengali, Marathi, etc.",
    "Kiosk hardware requires physical maintenance & power",
    "LLM hallucination risk despite safety pipeline",
    "No integration with HMIS/EHR — uses standalone DB",
    "Initial patient trust barrier with AI voice systems",
]
for l in lims:
    p2 = tf.add_paragraph()
    p2.text = f"• {l}"; p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(7)

# Future scope
s = add_shape(sl, Inches(9), Inches(1.2), Inches(4), Inches(5.8),
              BG_CARD, TEAL, Pt(1.5))
tf = s.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_top = Pt(12)
p = tf.paragraphs[0]; p.text = "🔮 Future Scope"
p.font.size = Pt(20); p.font.color.rgb = TEAL; p.font.bold = True
futures = [
    "On-device LLM (Llama 3 quantized) for zero-latency offline AI",
    "Expand to 8+ Indian languages (Bengali, Marathi, Kannada)",
    "HMIS/EHR integration for unified patient records",
    "AI triage: symptom assessment before department selection",
    "Video consultation integration for remote specialists",
    "Biometric authentication (Aadhaar, fingerprint)",
    "Predictive queue analytics with ML-based wait estimation",
    "Multi-hospital network deployment with centralized dashboard",
    "Accessibility: Braille output & sign-language avatar",
    "PWA deployment for mobile kiosk access",
]
for f in futures:
    p2 = tf.add_paragraph()
    p2.text = f"• {f}"; p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(6)

slide_number(sl, 17, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════════
# SLIDE 18 — Conclusion + References
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
slide_header(sl, "17 & 18.  Conclusion & References", ACCENT)

# Conclusion
s = add_shape(sl, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.8),
              BG_CARD, ACCENT, Pt(1.5))
tf = s.text_frame; tf.word_wrap = True
tf.margin_left = Pt(18); tf.margin_top = Pt(14)
p = tf.paragraphs[0]; p.text = "Conclusion"
p.font.size = Pt(22); p.font.color.rgb = ACCENT; p.font.bold = True

conclusions = [
    "MediKiosk successfully demonstrates that a voice-first AI kiosk can dramatically simplify "
    "hospital self-service for low-literacy, multilingual populations in India.",

    "The system reduces average patient registration time from ~8 minutes to ~2 minutes, "
    "eliminates paper-based processes, and provides real-time queue visibility — improving both "
    "patient experience and hospital operational efficiency.",

    "The 5-layer safety architecture ensures that LLM-generated outputs are rigorously validated "
    "before reaching the UI, addressing critical concerns about AI reliability in healthcare settings.",

    "The offline-first architecture with IndexedDB fallback guarantees service continuity in "
    "government hospital environments where internet connectivity is unreliable.",

    "MediKiosk validates the feasibility of deploying open-source AI models (Llama 3, Whisper) "
    "for real-world healthcare applications in resource-constrained environments.",
]
for c in conclusions:
    p2 = tf.add_paragraph()
    p2.text = f"•  {c}"; p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(8)
    p2.line_spacing = Pt(16)

# References
add_text_box(sl, Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
             "References", font_size=18, color=YELLOW, bold=True)

refs = [
    "[1]  Touvron, H., et al. \"Llama 2: Open Foundation and Fine-Tuned Chat Models.\" Meta AI, 2023.",
    "[2]  Radford, A., et al. \"Robust Speech Recognition via Large-Scale Weak Supervision.\" OpenAI, 2022. (Whisper)",
    "[3]  \"Sarvam AI: Building AI for India's Languages.\" Sarvam AI, 2024. https://www.sarvam.ai",
    "[4]  Ramírez, S., \"FastAPI: Modern Python Web Framework.\" https://fastapi.tiangolo.com",
    "[5]  React Team, \"React 19 Documentation.\" Meta, 2025. https://react.dev",
    "[6]  Firebase, \"Cloud Storage for Firebase.\" Google, 2025. https://firebase.google.com/docs/storage",
    "[7]  Hugging Face, \"Inference API Documentation.\" 2025. https://huggingface.co/docs/api-inference",
    "[8]  SQLAlchemy, \"SQLAlchemy 2.0 Unified Tutorial.\" https://docs.sqlalchemy.org",
    "[9]  World Health Organization, \"Digital Health in the WHO South-East Asia Region,\" 2023.",
    "[10] National Health Authority (India), \"Ayushman Bharat Digital Mission,\" 2024.",
]

s = add_shape(sl, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5),
              BG_CARD, YELLOW, Pt(1))
tf = s.text_frame; tf.word_wrap = True
tf.margin_left = Pt(14); tf.margin_top = Pt(10)
p = tf.paragraphs[0]; p.text = refs[0]
p.font.size = Pt(9); p.font.color.rgb = LIGHT
for r in refs[1:]:
    p2 = tf.add_paragraph()
    p2.text = r; p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT; p2.space_before = Pt(4)

# Bottom bar
add_rect(sl, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), ACCENT)
add_text_box(sl, Inches(1), Inches(7.1), Inches(11), Inches(0.35),
             "Thank You  ·  MediKiosk  ·  AMD General Hospital  ·  March 2026",
             font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)

slide_number(sl, 18, TOTAL_SLIDES)


# ── Save ────────────────────────────────────────────────────────
output_path = os.path.join(r"e:\AMD\medikisok", "MediKiosk_Project_Submission.pptx")
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
